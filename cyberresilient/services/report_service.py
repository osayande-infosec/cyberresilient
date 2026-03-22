"""
PDF Report Generator
Creates professional PDF reports with configurable branding.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from fpdf import FPDF

from cyberresilient.config import get_config

REPORTS_DIR = Path(__file__).resolve().parent.parent.parent / "reports"


def _ensure_reports_dir() -> None:
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)


class BrandedPDF(FPDF):
    """Custom PDF with configurable branding."""

    def header(self):
        cfg = get_config()
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(42, 42, 42)
        self.cell(0, 10, cfg.branding.app_title, align="L")
        self.set_font("Helvetica", "", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 10, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", align="R", new_x="LMARGIN", new_y="NEXT")
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(5)

    def footer(self):
        cfg = get_config()
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(150, 150, 150)
        self.cell(
            0, 10,
            f"{cfg.branding.app_title} — {cfg.branding.app_subtitle} | Page {self.page_no()}/{{nb}}",
            align="C",
        )

    def section_title(self, title: str):
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(42, 42, 42)
        self.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT")
        self.ln(2)

    def body_text(self, text: str):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(60, 60, 60)
        self.multi_cell(0, 6, text)
        self.ln(2)

    def kv_row(self, key: str, value: str, bold_value: bool = False):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(42, 42, 42)
        self.cell(60, 7, key)
        style = "B" if bold_value else ""
        self.set_font("Helvetica", style, 10)
        self.set_text_color(60, 60, 60)
        self.cell(0, 7, str(value), new_x="LMARGIN", new_y="NEXT")

    def status_badge(self, text: str, passed: bool):
        if passed:
            self.set_text_color(0, 128, 0)
        else:
            self.set_text_color(200, 0, 0)
        self.set_font("Helvetica", "B", 10)
        self.cell(0, 7, text, new_x="LMARGIN", new_y="NEXT")
        self.set_text_color(60, 60, 60)


def generate_dr_report(sim_result: dict, raci: list[dict]) -> str:
    """Generate a DR simulation test report PDF. Returns file path."""
    _ensure_reports_dir()
    pdf = BrandedPDF()
    pdf.alias_nb_pages()
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 18)
    pdf.set_text_color(42, 42, 42)
    pdf.cell(0, 15, "Disaster Recovery Simulation Report", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    pdf.section_title("1. Simulation Summary")
    pdf.kv_row("System:", sim_result["system_name"])
    pdf.kv_row("System ID:", sim_result["system_id"])
    pdf.kv_row("Scenario:", sim_result["scenario_name"])
    pdf.kv_row("Scenario Type:", sim_result["scenario_type"])
    pdf.kv_row("Severity:", sim_result["severity"])
    pdf.kv_row("DR Strategy:", sim_result["dr_strategy"])
    pdf.kv_row("Simulation Date:", sim_result["timestamp"])
    pdf.ln(3)

    pdf.section_title("2. RTO / RPO Analysis")
    rto_status = "PASS" if sim_result["rto_met"] else "FAIL"
    rpo_status = "PASS" if sim_result["rpo_met"] else "FAIL"
    overall = "PASS" if sim_result["overall_pass"] else "FAIL"

    pdf.kv_row("RTO Target:", f"{sim_result['rto_target_hours']} hours")
    pdf.kv_row("RTO Estimated:", f"{sim_result['rto_estimated_hours']} hours")
    pdf.status_badge(f"RTO Result: {rto_status}", sim_result["rto_met"])
    if sim_result["rto_gap_hours"] > 0:
        pdf.kv_row("RTO Gap:", f"{sim_result['rto_gap_hours']} hours")
    pdf.ln(2)

    pdf.kv_row("RPO Target:", f"{sim_result['rpo_target_hours']} hours")
    pdf.kv_row("RPO Estimated:", f"{sim_result['rpo_estimated_hours']} hours")
    pdf.status_badge(f"RPO Result: {rpo_status}", sim_result["rpo_met"])
    if sim_result["rpo_gap_hours"] > 0:
        pdf.kv_row("RPO Gap:", f"{sim_result['rpo_gap_hours']} hours")
    pdf.ln(2)

    pdf.status_badge(f"Overall Result: {overall}", sim_result["overall_pass"])
    pdf.ln(5)

    pdf.section_title("3. Impact Assessment")
    pdf.kv_row("Departments Affected:", ", ".join(sim_result["departments_affected"]))
    pdf.body_text(f"Public Impact: {sim_result['public_impact']}")
    pdf.ln(3)

    pdf.section_title("4. Recovery Steps")
    for i, step in enumerate(sim_result["recovery_steps"], 1):
        pdf.body_text(f"{i}. {step}")

    pdf.add_page()
    pdf.section_title("5. Recommendations")
    for rec in sim_result["recommendations"]:
        pdf.body_text(f"\u2022 {rec}")

    pdf.section_title("6. RACI Matrix \u2014 DR Activation")
    pdf.set_font("Helvetica", "B", 8)
    col_widths = [45, 30, 25, 45, 40]
    headers = ["Activity", "Responsible", "Accountable", "Consulted", "Informed"]
    for w, h in zip(col_widths, headers):
        pdf.cell(w, 7, h, border=1)
    pdf.ln()

    pdf.set_font("Helvetica", "", 7)
    for row in raci:
        vals = [row["activity"], row["responsible"], row["accountable"], row["consulted"], row["informed"]]
        for w, v in zip(col_widths, vals):
            pdf.cell(w, 7, v[:30], border=1)
        pdf.ln()

    filename = f"DR_Report_{sim_result['system_id']}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = REPORTS_DIR / filename
    pdf.output(str(filepath))
    return str(filepath)


def generate_risk_report(assessment: dict, vendor_name: str) -> str:
    """Generate an Architecture Risk Assessment PDF. Returns file path."""
    _ensure_reports_dir()
    pdf = BrandedPDF()
    pdf.alias_nb_pages()
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 15, "Architecture Security Assessment", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)

    pdf.section_title("1. Assessment Summary")
    pdf.kv_row("Vendor / Solution:", vendor_name)
    pdf.kv_row("Date:", datetime.now().strftime("%Y-%m-%d"))
    pdf.kv_row("Total Checks:", str(assessment["total_checks"]))
    pdf.kv_row("Passed:", str(assessment["passed"]))
    pdf.kv_row("Failed:", str(assessment["failed"]))
    pdf.kv_row("Score:", f"{assessment['score_pct']}%")
    pdf.status_badge(f"Overall Risk: {assessment['overall_risk']}", assessment["score_pct"] >= 70)
    pdf.ln(5)

    pdf.section_title("2. Detailed Findings")
    for r in assessment["results"]:
        status_str = "PASS" if r["passed"] else "FAIL"
        pdf.set_font("Helvetica", "B", 10)
        if r["passed"]:
            pdf.set_text_color(0, 128, 0)
        else:
            pdf.set_text_color(200, 0, 0)
        pdf.cell(0, 7, f"[{status_str}] {r['control']}", new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(60, 60, 60)
        pdf.set_font("Helvetica", "", 9)
        pdf.body_text(f"   Framework: {r['framework']}")
        if not r["passed"]:
            pdf.body_text(f"   Risk: {r['risk_if_missing']}")
            pdf.body_text(f"   Recommendation: {r['recommendation']}")
        pdf.ln(1)

    filename = f"ArchRisk_{vendor_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = REPORTS_DIR / filename
    pdf.output(str(filepath))
    return str(filepath)


def generate_executive_pptx(dashboard_data: dict) -> str:
    """Generate a PPTX executive security brief from dashboard data. Returns file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    _ensure_reports_dir()
    cfg = get_config()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
    GOLD = RGBColor(0xD4, 0xAF, 0x37)
    WHITE = RGBColor(0xEA, 0xEA, 0xEA)
    GREEN = RGBColor(0x4C, 0xAF, 0x50)
    RED = RGBColor(0xF4, 0x43, 0x36)

    def _set_slide_bg(slide, color=DARK_BG):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return tf

    # ── Slide 1: Title ──────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide1)
    _add_text(slide1, cfg.branding.app_title, 1, 1.5, 11, 1.5, font_size=40, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, "Executive Security Brief", 1, 3, 11, 1, font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, f"Generated: {datetime.now().strftime('%B %d, %Y')}", 1, 4.5, 11, 0.5, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, cfg.organization.name, 1, 5.2, 11, 0.5, font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Security Posture ───────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide2)
    _add_text(slide2, "Security Posture Overview", 0.5, 0.3, 12, 0.8, font_size=30, color=GOLD, bold=True)

    score = dashboard_data.get("overall_security_score", 0)
    score_color = GREEN if score >= 70 else RED
    _add_text(slide2, f"Overall Security Score: {score}/100", 0.5, 1.3, 6, 0.6, font_size=24, color=score_color, bold=True)

    kpis = dashboard_data.get("kpi_data", [])
    for i, kpi in enumerate(kpis[:8]):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 3.1
        y = 2.2 + row * 1.8
        label = kpi.get("label", "")
        val = kpi.get("value", "")
        unit = kpi.get("unit", "")
        _add_text(slide2, label, x, y, 2.8, 0.4, font_size=12, color=WHITE)
        _add_text(slide2, f"{val}{unit}", x, y + 0.4, 2.8, 0.5, font_size=22, color=GOLD, bold=True)

    # ── Slide 3: Risk Summary ──────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide3)
    _add_text(slide3, "Risk Landscape", 0.5, 0.3, 12, 0.8, font_size=30, color=GOLD, bold=True)

    from cyberresilient.services.risk_service import load_risks, get_risk_summary
    risks = load_risks()
    summary = get_risk_summary(risks)

    _add_text(slide3, f"Total Risks: {summary['total']}", 0.5, 1.3, 5, 0.5, font_size=20, color=WHITE)
    y_pos = 2.0
    level_colors = {"Very High": RED, "High": RGBColor(0xFF, 0x98, 0x00), "Medium": RGBColor(0xFF, 0xC1, 0x07), "Low": GREEN}
    for level, count in summary["by_level"].items():
        _add_text(slide3, f"{level}: {count}", 0.5, y_pos, 4, 0.4, font_size=16, color=level_colors.get(level, WHITE))
        y_pos += 0.5

    # Top risks
    _add_text(slide3, "Top Risks Requiring Attention", 6, 1.3, 6, 0.5, font_size=20, color=WHITE, bold=True)
    top_risks = sorted(risks, key=lambda r: r["risk_score"], reverse=True)[:5]
    y_pos = 2.0
    for r in top_risks:
        _add_text(slide3, f"• {r['title']} (Score: {r['risk_score']})", 6, y_pos, 6.5, 0.4, font_size=12, color=WHITE)
        y_pos += 0.45

    # ── Slide 4: Compliance ─────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide4)
    _add_text(slide4, "Compliance Status", 0.5, 0.3, 12, 0.8, font_size=30, color=GOLD, bold=True)

    compliance_trend = dashboard_data.get("compliance_trend", [])
    if compliance_trend:
        latest = compliance_trend[-1] if compliance_trend else {}
        nist_score = latest.get("nist_csf", 0)
        iso_score = latest.get("iso_27001", 0)
        _add_text(slide4, f"NIST CSF 2.0 Score: {nist_score}%", 0.5, 1.5, 5, 0.5, font_size=22, color=GREEN if nist_score >= 70 else RED, bold=True)
        _add_text(slide4, f"ISO 27001 Score: {iso_score}%", 0.5, 2.3, 5, 0.5, font_size=22, color=GREEN if iso_score >= 70 else RED, bold=True)

    threats = dashboard_data.get("threat_categories", [])
    if threats:
        _add_text(slide4, "Threat Landscape", 6, 1.3, 6, 0.5, font_size=20, color=WHITE, bold=True)
        y_pos = 2.0
        for t in threats[:6]:
            name = t.get("category", "")
            count = t.get("count", 0)
            _add_text(slide4, f"• {name}: {count} events", 6, y_pos, 6, 0.35, font_size=12, color=WHITE)
            y_pos += 0.4

    # ── Slide 5: Recommendations ────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide5)
    _add_text(slide5, "Key Recommendations", 0.5, 0.3, 12, 0.8, font_size=30, color=GOLD, bold=True)

    recommendations = [
        "Continue quarterly DR testing for all Tier 1 systems",
        "Address all 'Very High' risks within 30 days",
        "Complete NIST CSF gap remediation for Protect & Detect functions",
        "Expand security awareness training coverage to 100%",
        "Implement 24/7 SOC monitoring capability",
        "Review and update all expired security policies",
    ]
    for i, rec in enumerate(recommendations):
        _add_text(slide5, f"{i + 1}. {rec}", 0.5, 1.3 + i * 0.6, 12, 0.5, font_size=16, color=WHITE)

    _add_text(slide5, "Prepared by Security Operations  |  CONFIDENTIAL", 0.5, 6.5, 12, 0.4, font_size=10, color=GOLD, alignment=PP_ALIGN.CENTER)

    filename = f"Executive_Brief_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = REPORTS_DIR / filename
    prs.save(str(filepath))
    return str(filepath)
